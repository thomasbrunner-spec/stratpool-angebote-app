"""Extract DOCX/PDF/PPTX files to plain Markdown.

Used by the seeding pipeline to ingest historical offers and the Mini-Angebot
template. Run standalone:

    uv run --group seed python -m scripts.extract path/to/file.{docx,pdf,pptx}
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

import mammoth
from pptx import Presentation
from pypdf import PdfReader

_INLINE_IMAGE_RE = re.compile(r"!\[[^\]]*\]\([^)]*\)")
_MD_ESCAPE_RE = re.compile(r"\\([!\"#$%&'()*+,\-./:;<=>?@\[\\\]^_`{|}~])")


def docx_to_markdown(path: Path) -> str:
    with path.open("rb") as f:
        result = mammoth.convert_to_markdown(f)
    text = _INLINE_IMAGE_RE.sub("", result.value)
    text = _MD_ESCAPE_RE.sub(r"\1", text)
    return text.strip()


def pdf_to_markdown(path: Path) -> str:
    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append(f"## Page {i}\n\n{text}")
    return "\n\n".join(pages)


def pptx_to_markdown(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines: list[str] = [f"## Slide {i}"]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    slide_lines.append(text)
        if len(slide_lines) > 1:
            parts.append("\n\n".join(slide_lines))
    return "\n\n".join(parts)


_DISPATCH = {
    ".docx": docx_to_markdown,
    ".pdf": pdf_to_markdown,
    ".pptx": pptx_to_markdown,
}


def extract(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix not in _DISPATCH:
        raise ValueError(f"Unsupported format: {suffix} (file: {path.name})")
    return _DISPATCH[suffix](path)


def main() -> None:
    if len(sys.argv) != 2:
        print("Usage: python -m scripts.extract <file>", file=sys.stderr)
        sys.exit(2)
    path = Path(sys.argv[1])
    if not path.exists():
        print(f"File not found: {path}", file=sys.stderr)
        sys.exit(1)
    print(extract(path))


if __name__ == "__main__":
    main()
